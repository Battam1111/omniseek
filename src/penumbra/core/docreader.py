"""Document digestion — turn a document FILE (pptx/docx/xlsx/pdf/txt/md, local
path or URL) into readable, structured text + an HONEST media inventory.

The third perception primitive: web pages (penumbra_add_url), speech (penumbra_transcribe),
and now documents-as-files. Built P39 (2026-06-10) from a real 颜色框.pptx
(a 28-slide, 66-image figure-style reference deck): the eye digests uploads the
way it transcribes audio.

Engine choice (tested on THAT real deck, not on benchmarks): python-pptx vs
markitdown head-to-head — python-pptx keeps structure (per-slide outline, group
recursion, tables, notes, image→slide mapping) where markitdown returns one
markdown blob; office formats are zip+XML, so three small pure libs
(python-pptx / python-docx / openpyxl) beat one convert-everything dep. PDF
reuses the proven fitz path. All parser imports are lazy (asr.py pattern), so
importing this module costs nothing.

Honesty rule: much of a document's MEANING can live in its IMAGES (颜色框: 24 of
28 slides are pure pictures — its text layer is ~3 slides of block labels). Text
extraction alone would silently pretend the doc IS its text layer. So every
section carries a media count (the "where is the non-text meaning" map), and the
image half is delivered to the agent's own vision two ways:
  - read_document(export_media=True) writes images to disk (section-mapped names);
  - view_images() returns the actual image bytes for IN-BAND delivery (the server
    wraps them as MCP Image content) — no out-of-band scp. The cheap default is a
    CONTACT SHEET (downscaled thumbnails tiled into one labeled montage) so the
    agent triages 30 images for the cost of one, then pulls the few that matter at
    full res. The eye renders pixels; INTERPRETING them is the agent's vision —
    never a server-side captioner (a weaker model pre-chewing for the real one is
    the fabrication trap). eye = render + structure; agent = see + judge.

Transport: the eye runs on the host; the user's files live on client machines.
Convention: ``scp "<file>" <eye-host>:penumbra-inbox/`` (any window can),
then call with ``penumbra-inbox/<name>`` — relative paths resolve against the
service user's HOME. Exports land in ``penumbra-inbox/.exports/<stem>/``.
"""
from __future__ import annotations

import logging
import os
import re
import tempfile
from io import BytesIO
from pathlib import Path
from typing import Optional
from urllib.parse import unquote, urlparse

from penumbra.core import _netguard, _optdep, cache, safeurl

logger = logging.getLogger(__name__)


def _fitz():
    """PyMuPDF (fitz) is the optional [pdf] extra (AGPL); fail-open with an install hint."""
    return _optdep.require("fitz", "pdf")

# Code + config source files are just plain text, routed to the txt reader (no parser): a
# read returns the file's text + the honest "no extractable images" note. Lets penumbra_read_document /
# penumbra_add_url digest a raw repo file (roadmap-④ engineering-craft prereq).
_CODE_EXTS = ("py", "pyi", "ts", "tsx", "js", "jsx", "mjs", "cjs", "rs", "go", "java", "kt",
              "c", "h", "cpp", "cc", "hpp", "cs", "rb", "php", "swift", "scala", "lua", "r",
              "sh", "bash", "zsh", "sql", "toml", "cfg", "ini", "conf", "yaml", "yml",
              "json", "xml", "html", "css", "tex", "rst", "dockerfile", "makefile")
_SUPPORTED = ("pptx", "docx", "xlsx", "pdf", "txt", "md", "markdown", "csv") + _CODE_EXTS
_MAX_DOWNLOAD = 100 * 1024 * 1024   # 100MB URL-download cap
_MAX_TEXT_FILE = 2 * 1024 * 1024    # raw txt/md/csv read cap
_XLSX_MAX_ROWS = 300                # per sheet — honest truncation, not silence
_XLSX_MAX_COLS = 40
_MEDIA_LIST_CAP = 100               # inventory entries returned (export writes ALL)
_TTL = 3600                         # parsed-doc cache; local mtime+size in the key
_UA = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/121.0 Safari/537.36"

# Image-view (penumbra_view_doc_images) tuning.
_IMAGE_FORMATS = ("pptx", "pdf", "docx")  # formats that carry extractable embedded images
_VIEW_MAX_DIM = 1456                # full-res long-edge cap (Claude vision downsamples ~here anyway)
_VIEW_FULL_CAP = 12                 # max full-res images per call (beyond → contact sheet)
_SHEET_MAX = 30                     # max thumbnails per contact sheet (legible after downsample)
_VIEW_HARD_CAP = 300                # max images materialized per call (OOM guard for huge docs)
_SHEET_COLS = 4
_SHEET_CELL = 260                   # thumbnail box (px)
_SHEET_PAD = 8
_SHEET_LABEL = 18                   # caption strip height under each thumbnail (px)
# OCR (penumbra_read_document ocr=True) — text-in-pixels (scanned pages, baked-in labels, HEX codes).
_OCR_MAX_IMAGES = 60                # OCR is per-image work; cap per document for sanity
_OCR_MAX_DIM = 1800                 # downscale before OCR (speed; OCR gains nothing from huge images)
_ocr_engine = None                  # lazy RapidOCR singleton (load is expensive; keep warm)


def _clean_stem(filename: Optional[str]) -> str:
    """Embedded-image filename → a safe, bounded slug stem (shared by export + view)."""
    return re.sub(r"[^\w.-]+", "_", Path(filename or "").stem).strip("_")[:40] or "image"


def _walk(shapes):
    """Recurse into grouped pptx shapes so nested pictures / text are not missed."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk(sh.shapes)
        else:
            yield sh


def _fmt_of(src: str) -> str:
    """Format from the file extension of a path or URL (query string ignored)."""
    path = urlparse(src).path if "://" in src else src
    ext = Path(unquote(path)).suffix.lstrip(".").lower()
    return ext if ext in _SUPPORTED else ""


# Content-Type → our format vocab: types a downloaded resource by what the server ACTUALLY sends,
# the canonical way (an extension-less URL like arxiv.org/pdf/2203.02155v1 is served as
# application/pdf). The server's declared type is authoritative over a URL extension sniff.
_MIME_FMT = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "text/plain": "txt",
    "text/markdown": "md",
    "text/csv": "csv",
}


def _fmt_from_response(content_type: Optional[str], cd_filename: Optional[str]) -> str:
    """Resolve our format from a downloaded response's headers. A Content-Disposition filename
    extension is the most reliable signal (check it first); else map the Content-Type's MIME
    (stripping any ``; charset=…`` parameter). '' when neither is recognized."""
    if cd_filename:
        ext = Path(unquote(cd_filename)).suffix.lstrip(".").lower()
        if ext in _SUPPORTED:
            return ext
    mime = (content_type or "").split(";", 1)[0].strip().lower()
    return _MIME_FMT.get(mime, "")


def _allowed_doc_roots() -> list[Path]:
    """Directories penumbra_read_document / view_images may read from. Default: the inbox only.
    A deployer adds roots via PENUMBRA_DOC_ROOTS (',' or ':' separated absolute paths)."""
    roots = [(Path.home() / "penumbra-inbox").resolve()]
    for r in re.split(r"[:,]", os.environ.get("PENUMBRA_DOC_ROOTS", "")):
        r = r.strip()
        if r:
            try:
                roots.append(Path(os.path.expanduser(r)).resolve())
            except (OSError, ValueError):
                pass
    return roots


def _resolve_local(path: str) -> Path:
    """Local-path semantics: ~ expands; a relative path resolves against HOME (so the documented
    "penumbra-inbox/<name>" convention works verbatim). SANDBOXED: the resolved REAL path must sit
    under an allowed root (default ~/penumbra-inbox; extend via PENUMBRA_DOC_ROOTS), so a caller (or a
    prompt-injected agent) cannot read arbitrary host files like ~/.penumbra/credentials. `..` is
    collapsed by resolve() before the containment check. Raises PermissionError on escape."""
    p = Path(os.path.expanduser(path))
    if not p.is_absolute():
        p = Path.home() / p
    rp = p.resolve()
    roots = _allowed_doc_roots()
    if not any(rp == root or root in rp.parents for root in roots):
        raise PermissionError(
            "path is outside the allowed document roots "
            f"({', '.join(str(r) for r in roots)}); set PENUMBRA_DOC_ROOTS to permit more")
    return rp


def _window(text: str, start_char: int, max_chars: int) -> tuple[str, bool]:
    """(chunk, truncated) — mechanical windowing for big docs. When the cut lands mid-stream, snap
    back to the last newline within a bounded lookback so a page ends on a line boundary (not
    mid-row/mid-word); returned length stays EXACT so ``start_char += returned_chars`` paginates
    losslessly. No newline in range -> raw cut (never regresses, never empties)."""
    _SNAP_LOOKBACK = 400
    start = max(int(start_char or 0), 0)
    n = max(int(max_chars or 0), 1)
    end = start + n
    if end < len(text):                          # only snap when actually truncating
        lo = max(start + 1, end - _SNAP_LOOKBACK)  # start+1 clamp: chunk is NEVER empty (no paging stall)
        nl = text.rfind("\n", lo, end)
        if nl != -1:
            end = nl + 1                          # include the newline; the next page starts fresh
    chunk = text[start:end]
    return chunk, (start + len(chunk)) < len(text)


def _sec(label: str, lines: list[str], media: int) -> dict:
    body = "\n".join(ln for ln in lines if ln is not None)
    return {"label": label, "body": body.strip(), "media": media}


# ---------------------------------------------------------------------------
# Per-format readers — each returns (title, sections, media_records)
#   sections: [{label, body, media}]
#   media_records: [{section, name, bytes, px}] (+ blob writer via export_dir)
# ---------------------------------------------------------------------------

def _read_pptx(path: Path, export_dir: Optional[Path]) -> tuple[str, list[dict], list[dict]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    title = (prs.core_properties.title or "").strip() or path.name

    sections, media = [], []
    for i, slide in enumerate(prs.slides, 1):
        lines, n_pics = [], 0
        for sh in _walk(slide.shapes):
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                n_pics += 1
                rec = {"section": f"Slide {i}", "name": f"s{i:02d}_{n_pics:02d}",
                       "bytes": None, "px": None}
                try:
                    img = sh.image
                    rec["name"] = f"s{i:02d}_{n_pics:02d}_{_clean_stem(img.filename)}.{img.ext}"
                    rec["bytes"] = len(img.blob)
                    rec["px"] = list(img.size) if img.size else None
                    if export_dir is not None:
                        (export_dir / rec["name"]).write_bytes(img.blob)
                except Exception as exc:  # linked (not embedded) picture etc.
                    logger.debug("pptx image skip slide %d: %s", i, exc)
                media.append(rec)
                continue
            if getattr(sh, "has_table", False) and sh.has_table:
                for row in sh.table.rows:
                    lines.append(" | ".join((c.text or "").strip() for c in row.cells))
                continue
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).replace("\xa0", " ").strip()
                    if t:
                        lines.append(t)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"[speaker notes] {notes}")
        if n_pics:
            lines.append(f"[{n_pics} image(s) on this slide]")
        sections.append(_sec(f"Slide {i}", lines, n_pics))
    return title, sections, media


def _read_docx(path: Path, export_dir: Optional[Path]) -> tuple[str, list[dict], list[dict]]:
    import docx

    d = docx.Document(str(path))
    title = (d.core_properties.title or "").strip() or path.name
    lines: list[str] = []
    for para in d.paragraphs:
        t = para.text.replace("\xa0", " ").strip()
        if not t:
            continue
        m = re.match(r"heading\s*(\d)", (para.style.name or "").lower())
        lines.append(("#" * min(int(m.group(1)), 4) + " " + t) if m else t)
    for tb in d.tables:
        lines.append("[table]")
        for row in tb.rows:
            lines.append(" | ".join((c.text or "").strip() for c in row.cells))
    n_imgs = sum(1 for rel in d.part.rels.values() if "image" in rel.reltype)
    if n_imgs:
        lines.append(f"[{n_imgs} embedded image(s) in document]")
    return title, [_sec("Document", lines, n_imgs)], []


def _read_xlsx(path: Path, export_dir: Optional[Path]) -> tuple[str, list[dict], list[dict]]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sections = []
    for ws in wb.worksheets:
        lines, n = [], 0
        for row in ws.iter_rows(values_only=True):
            if n >= _XLSX_MAX_ROWS:
                lines.append(f"[...truncated at {_XLSX_MAX_ROWS} rows; sheet has more]")
                break
            cells = ["" if v is None else str(v) for v in row[:_XLSX_MAX_COLS]]
            if any(c.strip() for c in cells):
                lines.append("\t".join(cells).rstrip())
                n += 1
        sections.append(_sec(f"Sheet: {ws.title}", lines, 0))
    wb.close()
    return path.name, sections, []


def _read_pdf(path: Path, export_dir: Optional[Path]) -> tuple[str, list[dict], list[dict]]:
    fitz = _fitz()  # PyMuPDF — same engine as pdf_source (proven on real papers); optional [pdf] extra

    doc = fitz.open(str(path))
    title = (doc.metadata or {}).get("title") or path.name
    sections = []
    for i, page in enumerate(doc, 1):
        n_imgs = len(page.get_images(full=True))
        sections.append(_sec(f"Page {i}", [page.get_text("text").strip(),
                                           f"[{n_imgs} image(s) on this page]" if n_imgs else None],
                             n_imgs))
    doc.close()
    return title, sections, []


def _read_txt(path: Path, export_dir: Optional[Path]) -> tuple[str, list[dict], list[dict]]:
    raw = path.read_bytes()[:_MAX_TEXT_FILE]
    return path.name, [_sec(path.name, [raw.decode("utf-8", "replace")], 0)], []


_READERS = {"pptx": _read_pptx, "docx": _read_docx, "xlsx": _read_xlsx,
            "pdf": _read_pdf, "txt": _read_txt, "md": _read_txt,
            "markdown": _read_txt, "csv": _read_txt}
# Code/config files reuse the plain-text reader (code IS text — same windowing, no parser).
_READERS.update({ext: _read_txt for ext in _CODE_EXTS})


def _cd_filename(content_disposition: Optional[str]) -> Optional[str]:
    """Pull the filename from a Content-Disposition header (filename* or filename=)."""
    if not content_disposition:
        return None
    m = re.search(r"filename\*=(?:[^']*'[^']*')?([^;]+)", content_disposition, re.I)
    if m:
        return unquote(m.group(1).strip().strip('"'))
    m = re.search(r'filename="?([^";]+)"?', content_disposition, re.I)
    return m.group(1).strip() if m else None


def _download(url: str, fmt: str) -> tuple[Path, Optional[str], Optional[str]]:
    """Stream a URL to a temp file (bounded). Returns (path, content_type, cd_filename) so the
    caller can type a resource the URL extension couldn't (e.g. an extension-less arxiv.org/pdf/<id>
    served as application/pdf). ``fmt`` is only the temp-file suffix (".bin" when unknown)."""
    import httpx

    _blk = _netguard.security_block_reason(url)
    if _blk is not None:
        raise RuntimeError(f"refused SSRF-class url ({_blk}): {url[:120]}")
    # Follow redirects MANUALLY (follow_redirects=False) and re-validate EVERY hop's Location via
    # _netguard before connecting (C2 UNTRUSTED_URL close): this URL is agent-controlled and the bytes
    # come back in-band, so a blind 302 -> 169.254.169.254 was an SSRF oracle. safeurl.walk_redirects_
    # revalidated centralizes the per-hop guard (one _netguard decision, no forked SSRF logic) and hands
    # back the FINAL non-3xx response with its body unread so we still STREAM it under _MAX_DOWNLOAD.
    with httpx.Client(follow_redirects=False, timeout=90,
                      headers={"User-Agent": _UA}) as client:
        r = safeurl.walk_redirects_revalidated(client, "GET", url, max_redirects=10)
        try:
            r.raise_for_status()
            content_type = r.headers.get("content-type")
            cd_filename = _cd_filename(r.headers.get("content-disposition"))
            # mkstemp only AFTER the final response is validated + 2xx-checked: a refused redirect (raise
            # above) never leaves an empty temp file behind, unlike the old create-then-stream order.
            fd, tmp = tempfile.mkstemp(suffix=f".{fmt or 'bin'}", prefix="penumbra-doc-")
            os.close(fd)
            try:
                size = 0
                with open(tmp, "wb") as f:
                    for chunk in r.iter_bytes():
                        size += len(chunk)
                        if size > _MAX_DOWNLOAD:
                            raise RuntimeError(f"download exceeds {_MAX_DOWNLOAD // 2**20}MB cap")
                        f.write(chunk)
            except BaseException:  # noqa: BLE001 an oversize / stream failure must not orphan the temp file
                try:
                    os.remove(tmp)
                except OSError:
                    pass
                raise
        finally:
            r.close()
    return Path(tmp), content_type, cd_filename


def read_document(src: str, start_char: int = 0, max_chars: int = 24000,
                  export_media: bool = False, ocr: bool = False) -> dict:
    """Parse `src` (local path or URL) → structured readable text. See penumbra_read_document.

    ocr=True additionally runs OCR over every embedded image (pptx/pdf/docx) and folds the
    recognized text-in-pixels into the body under a '图中文字 (OCR)' section — turns a scanned
    doc / figure deck's baked-in text into searchable text. Slower (per-image); cached."""
    src = (src or "").strip()
    if not src:
        return {"source": src, "error": "empty path/url"}
    fmt = _fmt_of(src)
    is_url = "://" in src
    # A local path with no usable extension can only be typed by its name → reject before any IO.
    # A URL with no usable extension is NOT rejected here: the server's Content-Type / filename
    # types it after the fetch (the canonical way; an extension-less arxiv.org/pdf/<id> is a real
    # application/pdf). The unsupported error survives only for a genuinely unrecognized type.
    if not fmt and not is_url:
        return {"source": src,
                "error": f"unsupported or missing file extension (supported: {', '.join(_SUPPORTED)})"}

    tmp: Optional[Path] = None
    try:
        if is_url:
            ck_id = src
            tmp, content_type, cd_filename = _download(src, fmt)
            path = tmp
            if not fmt:
                fmt = _fmt_from_response(content_type, cd_filename)
                if not fmt:
                    return {"source": src, "error": (
                        f"unsupported document type: server sent content-type={content_type!r}"
                        f" (supported: {', '.join(_SUPPORTED)})")}
        else:
            try:
                path = _resolve_local(src)
            except PermissionError as exc:
                return {"source": src, "error": str(exc)}
            if not path.is_file():
                inbox = Path.home() / "penumbra-inbox"
                have = sorted(p.name for p in inbox.glob("*") if p.is_file())[:20] if inbox.is_dir() else []
                return {"source": src, "error": f"file not found: {path}",
                        "inbox_files": have}
            st = path.stat()
            ck_id = f"{path}|{st.st_mtime_ns}|{st.st_size}"

        export_dir: Optional[Path] = None
        if export_media:
            stem = re.sub(r"[^\w.-]+", "_", Path(unquote(urlparse(src).path if is_url else src)).stem) or "doc"
            export_dir = Path.home() / "penumbra-inbox" / ".exports" / stem
            export_dir.mkdir(parents=True, exist_ok=True)

        ck = cache.make_key("docreader", ck_id, fmt, bool(export_media), bool(ocr))
        rec = cache.get(ck)
        was_cached = rec is not None
        if rec is None:
            title, sections, media = _READERS[fmt](path, export_dir)
            full = "\n\n".join(f"## {s['label']}\n{s['body']}" if s["body"] else f"## {s['label']}"
                               for s in sections)
            rec = {
                "format": fmt, "title": title,
                "outline": [{"label": s["label"], "chars": len(s["body"]), "media": s["media"]}
                            for s in sections],
                "_full_text": full,
                "media_total": sum(s["media"] for s in sections),
                "media": media[:_MEDIA_LIST_CAP],
                "media_dir": str(export_dir) if export_dir else None,
            }
            if ocr and fmt in _IMAGE_FORMATS:
                blocks, n = [], 0
                for im in _iter_images(path, fmt, None, None):
                    if n >= _OCR_MAX_IMAGES:
                        blocks.append(f"[OCR 截断于 {_OCR_MAX_IMAGES} 张图]")
                        break
                    txt = ocr_image(im["data"])
                    if txt:
                        blocks.append(f"### {im['section_label']} · {im['name']}\n{txt}")
                        n += 1
                rec["ocr_images"] = n
                if blocks:
                    rec["_full_text"] += "\n\n## 图中文字 (OCR · 机器识别,可能有误)\n\n" + "\n\n".join(blocks)
            cache.set(ck, rec, ttl=_TTL)
        text, truncated = _window(rec["_full_text"], start_char, max_chars)
        out = {k: v for k, v in rec.items() if k != "_full_text"}
        out.update({"source": src, "text": text, "total_chars": len(rec["_full_text"]),
                    "returned_chars": len(text), "start_char": int(start_char or 0),
                    "truncated": truncated, "cached": was_cached})
        return out
    except Exception as exc:  # noqa: BLE001
        logger.warning("docreader failed %s: %s", src, exc)
        return {"source": src, "error": f"{type(exc).__name__}: {str(exc)[:300]}"}
    finally:
        if tmp is not None:
            try:
                os.remove(tmp)
            except OSError:
                pass


# ===========================================================================
# Image view (penumbra_view_doc_images) — deliver embedded images to the agent's
# vision IN-BAND. Tier 4 of the document modality: the eye renders pixels, the
# agent sees. Returns raw PNG bytes; server.py wraps them as MCP Image content.
# ===========================================================================

def _parse_sel(v, *, as_int: bool):
    """Comma/space-separated selection string (or list) → set, or None when empty.

    Robust across MCP clients: accepts "8,15, 25" / "8 15 25" / [8,15] for sections,
    or a list of image names. None/"" → None (meaning "no filter")."""
    if v is None:
        return None
    items = v if isinstance(v, (list, tuple)) else re.split(r"[,\s]+", str(v).strip())
    out = set()
    for it in items:
        s = str(it).strip()
        if not s:
            continue
        if as_int:
            try:
                out.add(int(s))
            except ValueError:
                continue
        else:
            out.add(s)
    return out or None


def _iter_pptx_images(path: Path, sel_sections, sel_names):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides, 1):
        if sel_sections and i not in sel_sections:
            continue
        n = 0
        for sh in _walk(slide.shapes):
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            n += 1
            try:
                img = sh.image
                name = f"s{i:02d}_{n:02d}_{_clean_stem(img.filename)}.{img.ext}"
                if sel_names and name not in sel_names:
                    continue
                yield {"section": i, "section_label": f"Slide {i}", "name": name,
                       "data": img.blob, "px": list(img.size) if img.size else None}
            except Exception as exc:  # linked (not embedded) picture
                logger.debug("pptx image skip slide %d #%d: %s", i, n, exc)


def _iter_pdf_images(path: Path, sel_sections, sel_names):
    fitz = _fitz()
    doc = fitz.open(str(path))
    try:
        for i, page in enumerate(doc, 1):
            if sel_sections and i not in sel_sections:
                continue
            n = 0
            for info in page.get_images(full=True):
                n += 1
                xref = info[0]
                try:
                    d = doc.extract_image(xref)
                    name = f"p{i:02d}_{n:02d}.{d.get('ext', 'png')}"
                    if sel_names and name not in sel_names:
                        continue
                    yield {"section": i, "section_label": f"Page {i}", "name": name,
                           "data": d["image"], "px": None}
                except Exception as exc:
                    logger.debug("pdf image skip page %d xref %s: %s", i, xref, exc)
    finally:
        doc.close()


def _iter_pdf_page_renders(path: Path, sel_pages, dpi: int = 150):
    """Render selected PDF PAGES to PNG (fitz get_pixmap) — the figure channel for VECTOR figures
    (TikZ/pgfplots architecture diagrams + training curves) that embed NO raster, so _iter_pdf_images
    reports media=0 and the figure is invisible. Same yield shape as the image iterators → flows
    through the same downscale + in-band delivery. Reuses the proven fitz engine; no new dependency."""
    fitz = _fitz()
    doc = fitz.open(str(path))
    try:
        for i, page in enumerate(doc, 1):
            if sel_pages and i not in sel_pages:
                continue
            try:
                pix = page.get_pixmap(dpi=dpi)
                yield {"section": i, "section_label": f"Page {i}", "name": f"page{i:02d}.png",
                       "data": pix.tobytes("png"), "px": f"{pix.width}x{pix.height}"}
            except Exception as exc:  # noqa: BLE001
                logger.debug("pdf page render skip %d: %s", i, exc)
    finally:
        doc.close()


def _iter_docx_images(path: Path, sel_sections, sel_names):
    import docx
    if sel_sections and 1 not in sel_sections:  # docx is one logical section
        return
    d = docx.Document(str(path))
    n = 0
    for rel in d.part.rels.values():
        if "image" not in rel.reltype:
            continue
        n += 1
        try:
            part = rel.target_part
            ext = (Path(str(part.partname)).suffix or ".png").lstrip(".")
            name = f"d_{n:02d}.{ext}"
            if sel_names and name not in sel_names:
                continue
            yield {"section": 1, "section_label": "Document", "name": name,
                   "data": part.blob, "px": None}
        except Exception as exc:
            logger.debug("docx image skip #%d: %s", n, exc)


def _iter_images(path: Path, fmt: str, sel_sections, sel_names):
    if fmt == "pptx":
        yield from _iter_pptx_images(path, sel_sections, sel_names)
    elif fmt == "pdf":
        yield from _iter_pdf_images(path, sel_sections, sel_names)
    elif fmt == "docx":
        yield from _iter_docx_images(path, sel_sections, sel_names)


def _pil_open_rgb(data: bytes):
    """Open bytes → a PIL image in a PNG/montage-safe mode (CMYK/P/LA normalized)."""
    from PIL import Image as PILImage
    im = PILImage.open(BytesIO(data))
    if im.mode not in ("RGB", "RGBA", "L"):
        im = im.convert("RGB")
    return im


def view_image_urls(urls, max_images: int = 8, max_dim: int = _VIEW_MAX_DIM) -> dict:
    """Download arbitrary image URLs → downscaled PNG bytes for IN-BAND delivery (server.py wraps
    them as MCP Image content). The URL counterpart of view_images (which works off a parsed
    document file): for SEEING the post images the eye surfaces in a doc's `media` field
    (xiaohongshu / zhihu / etc. — where the 干货 often lives) WITHOUT a download/scp dance.
    Browser-ish headers + a Referer so image CDNs (rednotecdn / xhscdn) serve us. Returns
    {"images": [{url, data:bytes, format} | {url, error}], "count": N_ok}."""
    import httpx
    items = urls if isinstance(urls, (list, tuple)) else re.split(r"[\s,]+", str(urls or "").strip())
    items = [u for u in (str(x).strip() for x in items) if u.startswith("http")][:max_images]
    if not items:
        return {"images": [], "count": 0, "error": "no image urls given"}
    hdrs = {"User-Agent": _UA, "Referer": "https://www.xiaohongshu.com/",
            "Accept": "image/avif,image/webp,image/*,*/*"}
    out: list[dict] = []
    for u in items:
        # SSRF guard: this fetch takes AGENT-controlled URLs and returns the bytes IN-BAND, so an
        # unguarded internal URL is an exfil oracle (and reaches the loopback CDP DevTools API on
        # 127.0.0.1:9222). The document-download path (this file) already guards; this one must too.
        _blk = _netguard.security_block_reason(u)
        if _blk:
            out.append({"url": u, "error": f"refused: {_blk}"})
            continue
        try:
            # Follow redirects MANUALLY + re-validate each hop via _netguard (C2 UNTRUSTED_URL close):
            # follow_redirects=True let a 302 -> 127.0.0.1:9222 reach the loopback CDP DevTools API and
            # exfil its bytes in-band. safeurl.walk_redirects_revalidated raises refused-SSRF on a blocked
            # hop (caught below -> error entry); the FINAL non-3xx response's bytes are used as before.
            with httpx.Client(follow_redirects=False, timeout=25, headers=hdrs) as client:
                r = safeurl.walk_redirects_revalidated(client, "GET", u, max_redirects=10)
                try:
                    r.read()  # buffer the final body (stream=True response) before touching .content
                    if r.status_code == 200 and len(r.content) > 500:
                        out.append({"url": u, "data": _downscale_png(r.content, max_dim), "format": "png"})
                    else:
                        out.append({"url": u, "error": f"HTTP {r.status_code} ({len(r.content)}B)"})
                finally:
                    r.close()
        except Exception as exc:  # noqa: BLE001
            out.append({"url": u, "error": f"{type(exc).__name__}: {str(exc)[:80]}"})
    return {"images": out, "count": sum(1 for o in out if o.get("data"))}


def _downscale_png(data: bytes, max_dim: int = _VIEW_MAX_DIM) -> bytes:
    """Re-encode to PNG, downscaled so the long edge ≤ max_dim (never upscale)."""
    from PIL import Image as PILImage
    im = _pil_open_rgb(data)
    if max(im.size) > max_dim:
        im.thumbnail((max_dim, max_dim), PILImage.LANCZOS)
    buf = BytesIO()
    im.save(buf, format="PNG")
    return buf.getvalue()


def _get_ocr():
    global _ocr_engine
    if _ocr_engine is None:
        # OCR is the optional [ocr] extra (rapidocr-onnxruntime / onnxruntime); fail-open with a hint.
        _ocr_engine = _optdep.require("rapidocr_onnxruntime", "ocr").RapidOCR()
    return _ocr_engine


def ocr_image(data: bytes) -> str:
    """Text-in-pixels → text (RapidOCR: CJK + Latin + digits, local / keyless / CPU). This is
    mechanical TRANSCRIPTION of printed text, NOT interpretation of the figure (that stays the
    agent's vision). Benchmarked on real palette cards: reads HEX/RGB + Chinese accurately at
    ~0.5s/image. Returns the recognized lines joined; '' when nothing legible."""
    from PIL import Image as PILImage
    import numpy as np
    im = _pil_open_rgb(data)
    if max(im.size) > _OCR_MAX_DIM:
        im.thumbnail((_OCR_MAX_DIM, _OCR_MAX_DIM), PILImage.LANCZOS)
    try:
        res, _ = _get_ocr()(np.array(im))
    except Exception as exc:  # noqa: BLE001
        logger.warning("ocr failed: %s", exc)
        return ""
    if not res:
        return ""
    return "\n".join(
        (line[1] or "").strip() for line in res
        if isinstance(line, (list, tuple)) and len(line) > 1 and (line[1] or "").strip()
    )


def _contact_sheet(cells: list[dict]) -> bytes:
    """Tile downscaled thumbnails into ONE labeled montage PNG for cheap triage.

    Each cell: the thumbnail centered in a box + a caption ("#idx slN") below it,
    so the agent can name which ones to pull full-res. ASCII labels → the always-
    present default bitmap font (no CJK-font dependency)."""
    from PIL import Image as PILImage, ImageDraw, ImageFont
    cols = max(1, min(_SHEET_COLS, len(cells)))
    rows = -(-len(cells) // cols)
    cell_h = _SHEET_CELL + _SHEET_LABEL
    width = _SHEET_PAD + cols * (_SHEET_CELL + _SHEET_PAD)
    height = _SHEET_PAD + rows * (cell_h + _SHEET_PAD)
    sheet = PILImage.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(sheet)
    font = ImageFont.load_default()
    for idx, c in enumerate(cells):
        r, cc = divmod(idx, cols)
        x = _SHEET_PAD + cc * (_SHEET_CELL + _SHEET_PAD)
        y = _SHEET_PAD + r * (cell_h + _SHEET_PAD)
        try:
            im = _pil_open_rgb(c["data"])
            im.thumbnail((_SHEET_CELL, _SHEET_CELL), PILImage.LANCZOS)
            sheet.paste(im, (x + (_SHEET_CELL - im.size[0]) // 2,
                             y + (_SHEET_CELL - im.size[1]) // 2))
        except Exception as exc:  # unreadable image → labeled empty cell
            logger.debug("contact-sheet thumb skip %s: %s", c.get("name"), exc)
            draw.rectangle([x, y, x + _SHEET_CELL, y + _SHEET_CELL], outline="lightgray")
        tag = c["section_label"].replace("Slide ", "sl").replace("Page ", "pg").replace("Document", "doc")
        draw.text((x + 2, y + _SHEET_CELL + 3), f"#{idx + 1} {tag}"[:30], fill="black", font=font)
    buf = BytesIO()
    sheet.save(buf, format="PNG")
    return buf.getvalue()


def view_images(src: str, sections=None, names=None, max_images: int = _VIEW_FULL_CAP,
                contact_sheet: bool = False, max_dim: int = _VIEW_MAX_DIM, render_pages=None) -> dict:
    """Extract embedded images from `src` for in-band delivery. See penumbra_view_doc_images.

    ``render_pages`` (PDF only): instead of extracting embedded rasters, RENDER those page numbers to
    PNG — the way to SEE vector/TikZ figures + architecture diagrams that carry no embedded image.

    Returns raw PNG bytes (server wraps as MCP Image):
      contact-sheet mode → {mode:"contact_sheet", sheet:<png>, manifest:[...], ...}
      full mode          → {mode:"full", images:[{name,section,data:<png>}], manifest, ...}
      or {note:...} when nothing matched / no images, or {error:...}.
    """
    src = (src or "").strip()
    if not src:
        return {"source": src, "error": "empty path/url"}
    fmt = _fmt_of(src)
    is_url = "://" in src
    # A local path with no usable extension is typed by its name only → reject before IO. A URL
    # with no usable extension defers to the server's Content-Type after download (canonical typing,
    # same helper as read_document); the _IMAGE_FORMATS gate is then re-checked on the resolved fmt.
    if not fmt and not is_url:
        return {"source": src,
                "error": f"unsupported or missing file extension (supported: {', '.join(_SUPPORTED)})"}
    if fmt and fmt not in _IMAGE_FORMATS:
        return {"source": src, "format": fmt, "total_images": 0, "images": [],
                "note": f"{fmt} carries no extractable embedded images "
                        f"(its meaning is text; use penumbra_read)"}

    tmp: Optional[Path] = None
    try:
        if is_url:
            tmp, content_type, cd_filename = _download(src, fmt)
            path = tmp
            if not fmt:
                fmt = _fmt_from_response(content_type, cd_filename)
                if not fmt:
                    return {"source": src, "error": (
                        f"unsupported document type: server sent content-type={content_type!r}"
                        f" (supported: {', '.join(_SUPPORTED)})")}
                if fmt not in _IMAGE_FORMATS:
                    return {"source": src, "format": fmt, "total_images": 0, "images": [],
                            "note": f"{fmt} carries no extractable embedded images "
                                    f"(its meaning is text; use penumbra_read)"}
        else:
            try:
                path = _resolve_local(src)
            except PermissionError as exc:
                return {"source": src, "error": str(exc)}
            if not path.is_file():
                inbox = Path.home() / "penumbra-inbox"
                have = sorted(p.name for p in inbox.glob("*") if p.is_file())[:20] if inbox.is_dir() else []
                return {"source": src, "error": f"file not found: {path}", "inbox_files": have}

        sel_sections = _parse_sel(sections, as_int=True)
        sel_names = _parse_sel(names, as_int=False)
        render_sel = _parse_sel(render_pages, as_int=True)
        # Materialize with a hard cap so a pathological doc (hundreds of images) can
        # never load unbounded bytes into memory; a real selection skips early anyway.
        items, capped = [], False
        src_iter = (_iter_pdf_page_renders(path, render_sel) if (render_sel and fmt == "pdf")
                    else _iter_images(path, fmt, sel_sections, sel_names))
        for it in src_iter:
            items.append(it)
            if len(items) >= _VIEW_HARD_CAP:
                capped = True
                break
        total = len(items)
        if total == 0:
            picked = bool(sel_sections or sel_names or render_sel)
            return {"source": src, "format": fmt, "total_images": 0, "images": [], "sheet": None,
                    "note": "no images matched that section/name selection" if picked
                            else "this document has no embedded images"}
        cap_note = (f" (capped at {_VIEW_HARD_CAP}; narrow with sections= to reach the rest)"
                    if capped else "")

        # Contact sheet when: explicitly asked, OR no selection (triage the whole
        # doc), OR the selection exceeds the full-res cap.
        as_sheet = contact_sheet or not (sel_sections or sel_names or render_sel) or total > max_images
        note = cap_note.strip() if cap_note else ""
        if as_sheet:
            cells = items[:_SHEET_MAX]
            if total > _SHEET_MAX:
                note = (f"contact sheet shows {_SHEET_MAX} of {total}{cap_note}; narrow with "
                        f"sections=/names= (see the penumbra_read outline) for the rest")
            elif total > max_images and not contact_sheet:
                note = (f"{total} images selected (> {max_images} full-res cap) → contact sheet; "
                        f"pass names=\"...\" to pull specific ones at full res")
            manifest = [{"idx": i + 1, "section": c["section"], "section_label": c["section_label"],
                         "name": c["name"]} for i, c in enumerate(cells)]
            return {"source": src, "format": fmt, "mode": "contact_sheet",
                    "sheet": _contact_sheet(cells), "manifest": manifest,
                    "total_images": total, "shown": len(cells), "note": note}

        chosen = items[:max_images]
        images = [{"section": c["section"], "section_label": c["section_label"], "name": c["name"],
                   "px": c["px"], "data": _downscale_png(c["data"], max_dim)} for c in chosen]
        manifest = [{"idx": i + 1, "section": o["section"], "section_label": o["section_label"],
                     "name": o["name"], "px": o["px"]} for i, o in enumerate(images)]
        return {"source": src, "format": fmt, "mode": "full", "images": images,
                "manifest": manifest, "total_images": total, "shown": len(images), "note": note}
    except Exception as exc:  # noqa: BLE001
        logger.warning("docreader view_images failed %s: %s", src, exc)
        return {"source": src, "error": f"{type(exc).__name__}: {str(exc)[:300]}"}
    finally:
        if tmp is not None:
            try:
                os.remove(tmp)
            except OSError:
                pass
